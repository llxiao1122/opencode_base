#!/usr/bin/env python3
"""
read_doc - 按需读取知识库文件内容（从二进制提取文本）
用法: python3 read_doc.py 相对路径
"""
import sys, json
from pathlib import Path

KB_DIR = Path("/home/admin/opencode/Knowledge")


def extract_docx(fpath: Path) -> str:
    from docx import Document
    doc = Document(str(fpath))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)


def extract_xlsx(fpath: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(fpath), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"=== {sheet} ===")
        for row in ws.iter_rows(values_only=True):
            vals = [str(v) if v is not None else "" for v in row]
            text = " | ".join(vals).strip()
            if text:
                lines.append(text)
    return "\n".join(lines)


def extract_pptx(fpath: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(fpath))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)


def extract_pdf(fpath: Path) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(str(fpath))
    lines = []
    for page in reader.pages:
        text = page.extract_text()
        if text.strip():
            lines.append(text.strip())
    return "\n".join(lines)


def extract_doc(fpath: Path) -> str:
    return f"[无法提取 .doc 格式（旧版Word），请打开原始文件查看]"


def extract_file(fpath: Path, max_chars=4000) -> dict:
    if not fpath.exists():
        return {"status": "error", "message": f"文件不存在: {fpath}"}

    ext = fpath.suffix.lower()
    try:
        if ext == ".docx":
            text = extract_docx(fpath)
        elif ext == ".xlsx":
            text = extract_xlsx(fpath)
        elif ext == ".pptx":
            text = extract_pptx(fpath)
        elif ext == ".pdf":
            text = extract_pdf(fpath)
        elif ext == ".doc":
            text = extract_doc(fpath)
        elif ext == ".docm":
            text = extract_docx(fpath)
        elif ext == ".xls":
            text = f"[无法提取 .xls 格式（旧版Excel），请打开原始文件查看]"
        elif ext == ".md":
            with open(fpath, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            text = f"[不支持的文件格式: {ext}]"

        if len(text) > max_chars:
            text = text[:max_chars] + f"\n\n... (已截断，全文 {len(text)} 字符)"

        return {"status": "success", "file": str(fpath), "content": text}
    except Exception as e:
        return {"status": "error", "file": str(fpath), "message": f"提取失败: {type(e).__name__}: {e}"}


def main():
    args = sys.argv[1:]
    if not args:
        print(json.dumps({"status": "error", "message": "请提供文件路径（相对于 Knowledge/）"}, ensure_ascii=False))
        sys.exit(1)

    rel_path = args[0]
    fpath = KB_DIR / rel_path

    result = extract_file(fpath)
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
